"""Abstract interface for pptx document loader implementations."""

import datetime
import logging
import mimetypes
import os
import re
import tempfile
import uuid
from urllib.parse import urlparse

import requests
from pptx import Presentation

from configs import dify_config
from core.rag.extractor.extractor_base import BaseExtractor
from core.rag.models.document import Document
from extensions.ext_database import db
from extensions.ext_storage import storage
from models.enums import CreatorUserRole
from models.model import UploadFile

logger = logging.getLogger(__name__)


class PPTXExtractor(BaseExtractor):
    """Load pptx files.

    Args:
        file_path: Path to the file to load.
    """

    def __init__(self, file_path: str, tenant_id: str, user_id: str):
        """Initialize with file path."""
        self.file_path = file_path
        self.tenant_id = tenant_id
        self.user_id = user_id

        if "~" in self.file_path:
            self.file_path = os.path.expanduser(self.file_path)

        # If the file is a web path, download it to a temporary file, and use that
        if not os.path.isfile(self.file_path) and self._is_valid_url(self.file_path):
            r = requests.get(self.file_path)

            if r.status_code != 200:
                raise ValueError(f"Check the url of your file; returned status code {r.status_code}")

            self.web_path = self.file_path
            # TODO: use a better way to handle the file
            self.temp_file = tempfile.NamedTemporaryFile()  # noqa SIM115
            self.temp_file.write(r.content)
            self.file_path = self.temp_file.name
        elif not os.path.isfile(self.file_path):
            raise ValueError(f"File path {self.file_path} is not a valid file or url")

    def __del__(self) -> None:
        if hasattr(self, "temp_file"):
            self.temp_file.close()

    def extract(self) -> list[Document]:
        """Load given path as single page."""
        content = self.parse_pptx(self.file_path)
        return [
            Document(
                page_content=content,
                metadata={"source": self.file_path},
            )
        ]

    @staticmethod
    def _is_valid_url(url: str) -> bool:
        """Check if the url is valid."""
        parsed = urlparse(url)
        return bool(parsed.netloc) and bool(parsed.scheme)

    def _extract_images_from_pptx(self, prs):
        image_map = {}
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_ext = shape.image.ext
                    if image_ext is None:
                        continue
                    # use uuid as file name
                    file_uuid = str(uuid.uuid4())
                    file_key = "image_files/" + self.tenant_id + "/" + file_uuid + "." + image_ext
                    mime_type, _ = mimetypes.guess_type(file_key)

                    storage.save(file_key, shape.image.blob)

                    # save file to db
                    upload_file = UploadFile(
                        tenant_id=self.tenant_id,
                        storage_type=dify_config.STORAGE_TYPE,
                        key=file_key,
                        name=file_key,
                        size=0,
                        extension=str(image_ext),
                        mime_type=mime_type or "",
                        created_by=self.user_id,
                        created_by_role=CreatorUserRole.ACCOUNT,
                        created_at=datetime.datetime.now(datetime.UTC).replace(tzinfo=None),
                        used=True,
                        used_by=self.user_id,
                        used_at=datetime.datetime.now(datetime.UTC).replace(tzinfo=None),
                    )

                    db.session.add(upload_file)
                    db.session.commit()
                    image_map[(slide_idx, shape.shape_id)] = (
                        f"![image]({dify_config.FILES_URL}/files/{upload_file.id}/file-preview)"
                    )

        return image_map

    def _table_to_markdown(self, table):
        markdown = []
        rows = table.rows
        cols = table.columns
        total_cols = len(cols)
        # Header
        header_cells = [cell.text.strip() for cell in rows[0].cells]
        markdown.append("| " + " | ".join(header_cells) + " |")
        markdown.append("| " + " | ".join(["---"] * total_cols) + " |")
        # Rows
        for row in list(rows)[1:]:
            row_cells = [cell.text.strip() for cell in row.cells]
            markdown.append("| " + " | ".join(row_cells) + " |")
        return "\n".join(markdown)

    def parse_pptx(self, pptx_path):
        prs = Presentation(pptx_path)

        content = []

        image_map = self._extract_images_from_pptx(prs)

        url_pattern = re.compile(r"http://[^\s+]+//|https://[^\s+]+")
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                # Extract text (including hyperlinks)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ""
                        for run in paragraph.runs:
                            run_text = run.text or ""
                            # Check for hyperlink
                            if run.hyperlink and run.hyperlink.address:
                                if url_pattern.match(run.hyperlink.address):
                                    para_text += f"[{run_text}]({run.hyperlink.address})"
                                else:
                                    para_text += run_text
                            else:
                                para_text += run_text
                        if para_text.strip():
                            slide_content.append(para_text.strip())
                # Extract images
                if hasattr(shape, "image"):
                    image_md = image_map.get((slide_idx, shape.shape_id))
                    if image_md:
                        slide_content.append(image_md)
                # Extract tables
                if shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    slide_content.append(table_md)
            if slide_content:
                content.append(f"# Slide {slide_idx + 1}\n" + "\n".join(slide_content))
        return "\n\n".join(content)
